# utils/file_parser.py
from PyPDF2 import PdfReader
import docx
import csv
import pptx

def parse_documents(uploaded_files):
    all_texts = []

    for file in uploaded_files:
        ext = file.name.split('.')[-1].lower()

        if ext == 'pdf':
            reader = PdfReader(file)
            all_texts.append(" ".join([p.extract_text() for p in reader.pages if p.extract_text()]))

        elif ext == 'docx':
            doc = docx.Document(file)
            all_texts.append(" ".join([p.text for p in doc.paragraphs]))

        elif ext == 'pptx':
            prs = pptx.Presentation(file)
            text = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text.append(shape.text)
            all_texts.append(" ".join(text))

        elif ext == 'csv':
            content = file.read().decode("utf-8").splitlines()
            reader = csv.reader(content)
            all_texts.append(" ".join([" ".join(row) for row in reader]))

        elif ext in ['txt', 'md']:
            text = file.read().decode("utf-8")
            all_texts.append(text)

    return all_texts
